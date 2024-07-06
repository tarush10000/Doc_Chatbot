import sys
import os
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QScrollBar,
                            QLineEdit, QPushButton, QListWidget, QListWidgetItem, QLabel, QTextEdit, QDialog)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from PyQt6.QtGui import QTextCursor
from docx import Document
from pptx import Presentation
import pdfplumber
from typing import List
import chromadb
import numpy as np
from sentence_transformers import CrossEncoder, SentenceTransformer
import google.generativeai as genai
import hashlib
import json
import edge_tts
import tempfile
import pygame
import asyncio
import shutil
import re

class CollapsibleReferenceWidget(QWidget):
    show_references_signal = pyqtSignal(str)

    def __init__(self, references):
        super().__init__()
        self.references = references
        self.is_collapsed = True
        self.init_ui()


    def init_ui(self):
        self.layout = QVBoxLayout()

        self.summary_label = QLabel(self.get_summary())
        self.summary_label.setWordWrap(True)
        self.summary_label.setStyleSheet("color: white;")

        self.toggle_button = QPushButton("Show More")
        self.toggle_button.setStyleSheet("background-color: #1E1E1E; color: #00A884; border: none;")
        self.toggle_button.clicked.connect(self.show_references)

        self.layout.addWidget(self.summary_label)
        self.layout.addWidget(self.toggle_button)

        self.setLayout(self.layout)

    def get_summary(self):
        return self.references.split('\n')[0]  # Display only the first line

    def show_references(self):
        self.show_references_signal.emit(self.references)

class ReferencesPopup(QDialog):
    def __init__(self, references):
        super().__init__()
        self.setWindowTitle("References")
        self.setModal(True)

        layout = QVBoxLayout()

        references_text = QTextEdit()
        references_text.setPlainText(references)
        references_text.setReadOnly(True)

        layout.addWidget(references_text)

        self.setLayout(layout)


class DocumentProcessor(QThread):
    finished = pyqtSignal(object)
    warning = pyqtSignal(str)

    def __init__(self, documents_folder, embeddings_folder, cache_file):
        super().__init__()
        self.documents_folder = documents_folder
        self.embeddings_folder = embeddings_folder
        self.cache_file = cache_file
        self.model = SentenceTransformer('all-MiniLM-L6-v2')

    def run(self):
        try:
            all_text = ""
            metadata = self.load_metadata()
            total_files = len([f for f in os.listdir(self.documents_folder) if f.endswith((".pdf", ".docx", ".pptx"))])
            processed_files = 0
            db = None
            files_in_folder = set(os.listdir(self.documents_folder))
            files_in_cache = set(metadata.keys())

            if not files_in_cache.intersection(files_in_folder):
                print("No files in the folder match the cache. Clearing cache and embeddings...")
                self.clear_cache_and_embeddings()
                metadata = {}

            for filename in files_in_folder:
                if filename.endswith((".pdf", ".docx", ".pptx")):
                    processed_files += 1
                    print(f"Processing file {processed_files} of {total_files}: {filename}")
                    file_path = os.path.join(self.documents_folder, filename)
                    file_checksum = self.compute_md5(file_path)

                    if filename in metadata and metadata[filename]['checksum'] == file_checksum:
                        print(f"Skipping {filename}, embeddings already exist.")
                        continue

                    print(f"Processing {filename}...")
                    try:
                        text = self.load_document(file_path)
                        all_text += text

                        chunked_text = self.split_text(text)
                        db = self.create_or_update_chroma_db(documents=chunked_text, filename=filename)

                        metadata[filename] = {'checksum': file_checksum}
                        self.save_metadata(metadata)
                        print(f"Completed processing {filename}")
                    except Exception as e:
                        warning_msg = f"Failed to process {filename}: {str(e)}"
                        print(warning_msg)
                        self.warning.emit(warning_msg)
                    print(f"Progress: {processed_files}/{total_files} files")

            print("All files processed. Embedding creation complete.")
            
            if db is None:
                print("No new documents were processed. Using existing database.")
                db = self.get_existing_db()
            
            self.finished.emit(db)
        except Exception as e:
            print(f"Error processing documents: {e}")

    def clear_cache_and_embeddings(self):
        # Clear cache.json
        if os.path.exists(self.cache_file):
            os.remove(self.cache_file)
            print("Cleared cache.json")

        # Clear embeddings_cache folder
        if os.path.exists(self.embeddings_folder):
            shutil.rmtree(self.embeddings_folder)
            os.makedirs(self.embeddings_folder)
            print("Cleared embeddings_cache folder")

    def load_metadata(self):
        if os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r') as f:
                    content = f.read()
                    if content.strip():  # Check if file is not empty
                        return json.loads(content)
                    else:
                        print("Cache file is empty. Starting with fresh metadata.")
                        return {}
            except json.JSONDecodeError:
                print("Invalid JSON in cache file. Starting with fresh metadata.")
                return {}
        else:
            print("Cache file does not exist. Starting with fresh metadata.")
            return {}

    def save_metadata(self, metadata):
        try:
            with open(self.cache_file, 'w') as f:
                json.dump(metadata, f, indent=4)
            print("Metadata saved successfully.")
        except Exception as e:
            print(f"Error saving metadata: {e}")

    def get_existing_db(self):
        chroma_client = chromadb.PersistentClient(path=self.embeddings_folder)
        try:
            return chroma_client.get_collection(name="rag_experiment")
        except ValueError:
            print("No existing database found. Creating a new empty database.")
            return chroma_client.create_collection(name="rag_experiment")

    def load_document(self, file_path):
        if file_path.endswith('.pdf'):
            return self.load_pdf(file_path)
        elif file_path.endswith('.docx'):
            return self.load_docx(file_path)
        elif file_path.endswith('.pptx'):
            return self.load_pptx(file_path)
        else:
            print(f"Unsupported file format: {file_path}")
            return ""

    def load_pdf(self, file_path):
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error loading PDF {file_path}: {e}")
        return text

    def load_docx(self, file_path):
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            print(f"Error loading DOCX {file_path}: {e}")
        return text

    def load_pptx(self, file_path):
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error loading PPTX {file_path}: {e}")
        return text

    def split_text(self, text: str) -> List[str]:
        chunks = []
        current_chunk = ""
        in_code_block = False
        
        for line in text.split('\n'):
            if line.strip().startswith('```'):
                in_code_block = not in_code_block
                current_chunk += line + '\n'
            elif in_code_block:
                current_chunk += line + '\n'
            else:
                if len(current_chunk) + len(line) > 1000:
                    chunks.append(current_chunk.strip())
                    current_chunk = line + '\n'
                else:
                    current_chunk += line + '\n'
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks

    def create_or_update_chroma_db(self, documents: List[str], filename: str):
        chroma_client = chromadb.PersistentClient(path=self.embeddings_folder)
        
        try:
            db = chroma_client.get_collection(name="rag_experiment")
            print(f"Updating existing collection for {filename}")
        except ValueError:
            db = chroma_client.create_collection(name="rag_experiment")
            print(f"Created new collection for {filename}")

        existing_ids = set(db.get(include=['documents'])['ids'])
        new_documents = []
        new_ids = []

        for i, d in enumerate(documents):
            doc_id = f"{filename}_{i}"
            if doc_id not in existing_ids:
                new_documents.append(d)
                new_ids.append(doc_id)

        if new_documents:
            db.add(documents=new_documents, ids=new_ids)
            print(f"Added {len(new_documents)} new chunks for {filename}")
        else:
            print(f"No new chunks to add for {filename}")

        return db

    def load_metadata(self):
        if os.path.exists(self.cache_file):
            try:
                with open(self.cache_file, 'r') as f:
                    content = f.read()
                    if content.strip():  # Check if file is not empty
                        return json.loads(content)
                    else:
                        print("Cache file is empty. Starting with fresh metadata.")
                        return {}
            except json.JSONDecodeError:
                print("Invalid JSON in cache file. Starting with fresh metadata.")
                return {}
        else:
            print("Cache file does not exist. Starting with fresh metadata.")
            return {}

    def save_metadata(self, metadata):
        try:
            with open(self.cache_file, 'w') as f:
                json.dump(metadata, f, indent=4)
            print("Metadata saved successfully.")
        except Exception as e:
            print(f"Error saving metadata: {e}")

    def compute_md5(self, file_path):
        md5_hash = hashlib.md5()
        with open(file_path, 'rb') as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                md5_hash.update(byte_block)
        return md5_hash.hexdigest()

class ChatBot(QThread):
    response_ready = pyqtSignal(str, str)

    def __init__(self, db, query, full_history, summary_history):
        super().__init__()
        self.db = db
        self.query = query
        self.full_history = full_history
        self.summary_history = summary_history

    def run(self):
        relevant_passages = self.get_relevant_passages(self.query, self.db, n_results=10)
        if not relevant_passages:
            answer = self.generate_answer(self.query)
            self.response_ready.emit(answer, "No specific references")
            return

        cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
        ranked_docs = self.rank_documents(cross_encoder, self.query, relevant_passages)
        if not ranked_docs:
            answer = self.generate_answer(self.query)
            self.response_ready.emit(answer, "No specific references")
            return

        prompt = self.make_rag_prompt(self.query, list(ranked_docs.values()))
        answer = self.generate_answer(prompt)
        if not answer.strip():
            answer = self.generate_answer(self.query)

        references = "*References: " + ", ".join(list(ranked_docs.keys())[:3]) + "*"
        self.response_ready.emit(answer, references)

    def get_relevant_passages(self, query: str, db, n_results: int) -> List[str]:
        try:
            results = db.query(query_texts=[query], n_results=n_results)['documents']
            return [item for sublist in results for item in sublist]
        except Exception as e:
            print(f"Error querying relevant passages: {e}")
            return []

    def rank_documents(self, cross_encoder: CrossEncoder, query: str, retrieved_documents: List[str]):
        try:
            pairs = [[query, doc] for doc in retrieved_documents]
            scores = cross_encoder.predict(pairs)
            ranks = np.argsort(scores)[::-1]
            ranked_docs = {os.path.basename(doc): doc for doc, rank_num in zip(retrieved_documents, ranks)}
            return ranked_docs
        except Exception as e:
            print(f"Error ranking documents: {e}")
            return {}

    def make_rag_prompt(self, query: str, relevant_passages: List[str]) -> str:
        escaped = " ".join(relevant_passages).replace("'", "").replace('"', "").replace("\n", " ")
        history_text = " ".join([f"{sender}: {msg}" for sender, msg in self.full_history])
        summary_text = " ".join(self.summary_history)
        prompt = f"""You are Doccy, a helpful and knowledgeable AI teacher with a charming personality who loves to solve students' doubts. 
        Use the information from the reference passages included below to answer the student's question. 
        Provide a clear, detailed, and to-the-point explanation of the requested topic, include relevant examples if the user asks for it.
        If the passage is irrelevant to the answer, you may ignore it.
        QUESTION: '{query}'
        FULL HISTORY: '{history_text}'
        SUMMARY HISTORY: '{summary_text}'
        PASSAGES: '{escaped}'
        ANSWER:
        """
        return prompt


    def generate_answer(self, prompt: str) -> str:
        gemini_api_key = os.getenv("Gemini_API")
        if not gemini_api_key:
            raise ValueError("Gemini API Key not provided. Please provide Gemini_API as an environment variable")
        genai.configure(api_key=gemini_api_key)
        model = genai.GenerativeModel('gemini-1.5-pro')
        try:
            answer = model.generate_content(prompt)
            return answer.text
        except Exception as e:
            print(f"Error generating answer: {e}")
            return ""

class TextToSpeech(QThread):
    finished = pyqtSignal()

    def __init__(self, text):
        super().__init__()
        self.text = text
        self.is_playing = False
        self.temp_filename = ""

    def preprocess_text(self, text):
        # Remove emojis and asterisks
        text = re.sub(r'[^\w\s,]', '', text)  # Remove non-word characters except spaces and commas
        text = re.sub(r'\*', '', text)        # Remove asterisks
        return text

    def run(self):
        asyncio.run(self.tts_and_play())

    async def tts_and_play(self):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            self.temp_filename = temp_file.name
        
        communicate = edge_tts.Communicate(self.preprocess_text(self.text), "en-US-EmmaMultilingualNeural")
        await communicate.save(self.temp_filename)
        
        try:
            pygame.mixer.init()
            pygame.mixer.music.load(self.temp_filename)
            pygame.mixer.music.play()
            self.is_playing = True
            
            while pygame.mixer.music.get_busy() and self.is_playing:
                pygame.time.Clock().tick(10)
            
            pygame.mixer.quit()
        except pygame.error as e:
            print(f"Error playing audio: {e}")
        finally:
            self.cleanup()
        self.finished.emit()

    def stop(self):
        self.is_playing = False
        # Handle pygame mixer within the same thread context
        try:
            pygame.mixer.music.stop()
        except pygame.error:
            pass
        self.cleanup()

    def cleanup(self):
        if self.temp_filename != "" and os.path.exists(self.temp_filename):
            try:
                os.unlink(self.temp_filename)
            except PermissionError:
                print(f"PermissionError: Unable to delete file {self.temp_filename}")
            self.temp_filename = ""

class ChatInterface(QWidget):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.full_history = []  # To store full chat history
        self.summary_history = []  # To store summarized chat history

    def init_ui(self):
        layout = QVBoxLayout()

        # Chat display
        self.chat_list = QListWidget()
        self.chat_list.setStyleSheet("""
            QListWidget {
                background-color: #121B22;
                border: none;
            }
            QListWidget::item {
                border: none;
                margin: 5px;
            }
            QScrollBar:vertical {
                background-color: #121B22;
                width: 12px;
                margin: 16px 0 16px 0;
                border: 1px solid #121B22;
            }
            QScrollBar::handle:vertical {
                background-color: #888888;
                min-height: 20px;
                border-radius: 6px;
            }
            QScrollBar::handle:vertical:hover {
                background-color: #555555;
            }
            QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {
                border: none;
                background: none;
            }
            QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical {
                background: none;
            }
        """)
        layout.addWidget(self.chat_list)

        # Input area
        chat_input_layout = QHBoxLayout()
        self.chat_input = QLineEdit()
        self.chat_input.setStyleSheet("""
            background-color: #2A2F32;
            color: white;
            border: 1px solid #3A3F41;
            border-radius: 20px;
            padding: 10px 15px;
            font-size: 14px;
        """)
        self.chat_input.setPlaceholderText("Type a message...")
        self.send_button = QPushButton('Send')
        self.send_button.setStyleSheet("background-color: #00A884; color: white;")
        self.stop_tts_button = QPushButton('Stop TTS')
        self.stop_tts_button.setStyleSheet("background-color: #FF4136; color: white;")
        self.stop_tts_button.setEnabled(False)
        chat_input_layout.addWidget(self.chat_input)
        chat_input_layout.addWidget(self.send_button)
        chat_input_layout.addWidget(self.stop_tts_button)

        layout.addLayout(chat_input_layout)
        self.setLayout(layout)

    def add_message(self, sender, message, is_warning=False):
        self.full_history.append((sender, message))
        summary = f"{sender}: {message}"
        self.summary_history.append(summary)

        item = QListWidgetItem()
        widget = QLabel(message)
        widget.setWordWrap(True)

        bg_color = '#1E1E1E'  # Default background color
        
        item = QListWidgetItem()
        if sender == 'DoccyReference':
            widget = CollapsibleReferenceWidget(message)
        else:
            widget = QLabel(message)
            widget.setWordWrap(True)
            if is_warning:
                bg_color = '#FF4136'  # Red background for warnings
            else:
                if sender == 'Doccy':
                    bg_color = '#6495ED'
                elif sender == 'User':
                    bg_color = '#005C4B'
                else:
                    bg_color = '#1E1E1E'

        widget.setStyleSheet(f"""
            background-color: {bg_color};
            color: white;
            border-radius: 10px;
            padding: 15px;
            max-width: 700px;
        """)
        widget.adjustSize()
        size = widget.sizeHint()
        size.setHeight(size.height() + 5)  # Add extra height to ensure text is not cut off
        item.setSizeHint(size)

        self.chat_list.addItem(item)
        self.chat_list.setItemWidget(item, widget)

        if sender == 'User':
            item.setTextAlignment(Qt.AlignmentFlag.AlignRight)
            widget.setStyleSheet(widget.styleSheet() + "margin-left: 180px;")
        else:
            item.setTextAlignment(Qt.AlignmentFlag.AlignLeft)
            widget.setStyleSheet(widget.styleSheet() + "margin-right: 180px;")

        self.chat_list.scrollToBottom()

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.tts = None

    def init_ui(self):
        self.setWindowTitle("Doccy - Your Document Chatbot")
        self.setGeometry(100, 100, 800, 600)
        self.setStyleSheet("background-color: #121B22;")

        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        layout = QVBoxLayout()
        self.chat_interface = ChatInterface()
        layout.addWidget(self.chat_interface)

        central_widget.setLayout(layout)

        # Initialize document processor
        self.doc_processor = DocumentProcessor("documents", "embeddings_cache", "cache.json")
        self.doc_processor.finished.connect(self.on_processing_finished)
        self.doc_processor.warning.connect(self.show_warning)
        self.doc_processor.start()

    def on_processing_finished(self, db):
        self.db = db
        # Connect send button and chat input
        self.chat_interface.send_button.clicked.connect(self.send_message)
        self.chat_interface.chat_input.returnPressed.connect(self.send_message)
        self.chat_interface.stop_tts_button.clicked.connect(self.stop_tts)

    def show_warning(self, message):
        self.chat_interface.add_message('System', message, is_warning=True)

    def send_message(self):
        user_message = self.chat_interface.chat_input.text()
        if user_message.strip() == "":
            return

        self.chat_interface.add_message('User', user_message)
        self.chat_interface.chat_input.clear()

        self.chat_interface.add_message('Bot', '...')
        QApplication.processEvents()

        try:
            self.chatbot = ChatBot(self.db, user_message, self.chat_interface.full_history, self.chat_interface.summary_history)
            self.chatbot.response_ready.connect(self.on_response_ready)
            self.chatbot.start()
        except Exception as e:
            self.chat_interface.chat_list.takeItem(self.chat_interface.chat_list.count() - 1)
            error_message = f"An error occurred: {str(e)}. Please try again later or check your internet connection."
            self.chat_interface.add_message('Bot', error_message, is_warning=True)

    def show_references_popup(self, references):
        popup = ReferencesPopup(references)
        popup.exec()
    
    def on_response_ready(self, response, references):
        self.chat_interface.chat_list.takeItem(self.chat_interface.chat_list.count() - 1)
        self.chat_interface.add_message('Doccy', response)

        # Add references as collapsible widget
        references_widget = CollapsibleReferenceWidget(references)
        references_widget.show_references_signal.connect(self.show_references_popup)

        item = QListWidgetItem()
        item.setSizeHint(references_widget.sizeHint())
        self.chat_interface.chat_list.addItem(item)
        self.chat_interface.chat_list.setItemWidget(item, references_widget)

        # Stop previous TTS if it's still running
        self.stop_tts()

        # Start new text-to-speech
        self.tts = TextToSpeech(response)
        self.tts.start()
        self.chat_interface.stop_tts_button.setEnabled(True)

    def stop_tts(self):
        if self.tts:
            self.tts.stop()
            self.tts = None
        self.chat_interface.stop_tts_button.setEnabled(False)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())